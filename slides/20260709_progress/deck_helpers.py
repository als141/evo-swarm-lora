# -*- coding: utf-8 -*-
"""python-pptx 用のスライド組み立てヘルパー。源暎ゴシックP・フラット学術デザイン。"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from copy import deepcopy

EMU_IN = 914400

# ---- カラーパレット（AIっぽくないインク基調）----
INK    = RGBColor(0x22, 0x30, 0x3C)   # 濃紺グレー
SUB    = RGBColor(0x5A, 0x6B, 0x78)   # 補助
FAINT  = RGBColor(0x93, 0x9E, 0xA8)
PAPER  = RGBColor(0xFB, 0xFA, 0xF7)   # 生成りの紙
CARD   = RGBColor(0xFF, 0xFF, 0xFF)
LINE   = RGBColor(0xD8, 0xDE, 0xE3)
ACCENT = RGBColor(0xC6, 0x7A, 0x4E)   # くすんだ橙（アクセント）
TEAL   = RGBColor(0x2E, 0x6E, 0x8E)   # 主役の青
POS    = RGBColor(0x3F, 0x7D, 0x5A)   # 深緑
NEG    = RGBColor(0xB4, 0x55, 0x3F)   # レンガ
GOLD   = RGBColor(0xB8, 0x8A, 0x3D)

CRITIC  = RGBColor(0x3B, 0x6E, 0xA5)  # 批判＝青
PRAG    = RGBColor(0x4E, 0x8D, 0x6E)  # 実務＝緑
EXPLORE = RGBColor(0x8E, 0x6B, 0xB0)  # 探索＝紫
CRITIC_BG  = RGBColor(0xEC, 0xF1, 0xF7)
PRAG_BG    = RGBColor(0xEA, 0xF3, 0xEE)
EXPLORE_BG = RGBColor(0xF1, 0xEC, 0xF6)
INK_PANEL  = RGBColor(0x1C, 0x28, 0x33)

FONT = "GenEi Gothic P"
FONT_L = "GenEi Gothic P Light"
FONT_M = "GenEi Gothic P Medium"
FONT_SB = "GenEi Gothic P SemiBold"

W, H = Inches(13.333), Inches(7.5)


def new_deck():
    p = Presentation()
    p.slide_width = W
    p.slide_height = H
    return p


def _set_run_font(run, name):
    """latin/ea/cs すべてに typeface を設定して日本語フォールバックを防ぐ。"""
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def add_slide(prs, bg=PAPER):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    if bg is not None:
        r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
        r.fill.solid(); r.fill.fore_color.rgb = bg
        r.line.fill.background()
        r.shadow.inherit = False
        # 背景を最背面へ
        sp = r._element
        sp.getparent().remove(sp)
        s.shapes._spTree.insert(2, sp)
    return s


def rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0,
         shape=MSO_SHAPE.RECTANGLE, radius=0.10, shadow=False):
    sh = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    if shadow:
        _soft_shadow(sh)
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sh.adjustments[0] = radius
        except Exception:
            pass
    return sh


def _soft_shadow(shape):
    spPr = shape._element.spPr
    lst = spPr.makeelement(qn("a:effectLst"), {})
    sh = lst.makeelement(qn("a:outerShdw"),
                         {"blurRad": "60000", "dist": "25000", "dir": "5400000",
                          "rotWithShape": "0"})
    clr = sh.makeelement(qn("a:srgbClr"), {"val": "222833"})
    alpha = clr.makeelement(qn("a:alpha"), {"val": "16000"})
    clr.append(alpha); sh.append(clr); lst.append(sh); spPr.append(lst)


def line(slide, x, y, w, h, color=LINE, weight=1.2):
    ln = slide.shapes.add_connector(2, Inches(x), Inches(y), Inches(x+w), Inches(y+h))
    ln.line.color.rgb = color; ln.line.width = Pt(weight)
    ln.shadow.inherit = False
    return ln


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         line_spacing=1.12, wrap=True, space_after=2):
    """runs: [(text, size, color, font, {opts})] または文字列。段落は \n で分割。"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    if isinstance(runs, str):
        runs = [(runs, 18, INK, FONT, {})]
    # 段落単位: runs はフラットに与え、seg text 内 \n で改行
    first = True
    para = tf.paragraphs[0]
    para.alignment = align
    if line_spacing:
        para.line_spacing = line_spacing
    para.space_after = Pt(space_after); para.space_before = Pt(0)
    for seg in runs:
        txt, size, color, fnt, opts = _norm(seg)
        lines = txt.split("\n")
        for li, part in enumerate(lines):
            if li > 0:
                para = tf.add_paragraph()
                para.alignment = align
                if line_spacing:
                    para.line_spacing = line_spacing
                para.space_after = Pt(space_after); para.space_before = Pt(0)
            r = para.add_run(); r.text = part
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = opts.get("bold", False)
            if opts.get("italic"):
                r.font.italic = True
            _set_run_font(r, fnt)
    return tb


def _norm(seg):
    if isinstance(seg, str):
        return seg, 18, INK, FONT, {}
    txt = seg[0]
    size = seg[1] if len(seg) > 1 else 18
    color = seg[2] if len(seg) > 2 else INK
    fnt = seg[3] if len(seg) > 3 else FONT
    opts = seg[4] if len(seg) > 4 else {}
    return txt, size, color, fnt, opts


def para_block(slide, x, y, w, h, paragraphs, align=PP_ALIGN.LEFT,
               anchor=MSO_ANCHOR.TOP, line_spacing=1.22, space_after=8):
    """paragraphs: [[(text,size,color,font,opts), ...], ...] 各要素が1段落。"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for m in (tf.margin_left, tf.margin_right, tf.margin_top, tf.margin_bottom):
        pass
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, pg in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = line_spacing
        p.space_after = Pt(space_after); p.space_before = Pt(0)
        for seg in pg:
            txt, size, color, fnt, opts = _norm(seg)
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = opts.get("bold", False)
            if opts.get("italic"):
                r.font.italic = True
            _set_run_font(r, fnt)
    return tb


def bullets(slide, x, y, w, h, items, size=17, gap=10, color=INK,
            marker="—", marker_color=ACCENT, line_spacing=1.2):
    """items: [str | (str, opts)]。markerで箇条書き。"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, it in enumerate(items):
        if isinstance(it, tuple):
            txt, opts = it
        else:
            txt, opts = it, {}
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(gap); p.space_before = Pt(0)
        rm = p.add_run(); rm.text = marker + "  "
        rm.font.size = Pt(size); rm.font.color.rgb = opts.get("marker_color", marker_color)
        _set_run_font(rm, FONT_SB)
        r = p.add_run(); r.text = txt
        r.font.size = Pt(opts.get("size", size))
        r.font.color.rgb = opts.get("color", color)
        r.font.bold = opts.get("bold", False)
        _set_run_font(r, opts.get("font", FONT))
    return tb


def title_head(slide, kicker, title, x=0.85, y=0.62):
    """本編スライドの見出し（小見出し＋大見出し＋アクセント線）。"""
    if kicker:
        text(slide, x, y, 11.6, 0.4, [(kicker, 13, ACCENT, FONT_SB, {})])
    text(slide, x, y + (0.34 if kicker else 0.0), 11.8, 0.8,
         [(title, 27, INK, FONT, {"bold": True})])
    line(slide, x, y + (0.34 if kicker else 0.0) + 0.72, 1.1, 0, color=ACCENT, weight=2.6)


def footer(slide, page, total, short="LoRAエージェント集団の進化 ／ 研究進捗"):
    text(slide, 0.85, 7.06, 8.0, 0.3, [(short, 9.5, FAINT, FONT_L, {})],
         anchor=MSO_ANCHOR.MIDDLE)
    text(slide, 11.6, 7.06, 0.9, 0.3, [(f"{page} / {total}", 9.5, FAINT, FONT_L, {})],
         align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def chip(slide, x, y, w, label, fill, txt_color=CARD, size=11, h=0.34):
    sh = rect(slide, x, y, w, h, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    tf = sh.text_frame; tf.word_wrap = False
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.size = Pt(size); r.font.color.rgb = txt_color; r.font.bold = True
    _set_run_font(r, FONT_SB)
    sh.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return sh


def pic_fit(slide, path, x, y, w, h, align="center", valign="middle"):
    """アスペクト比を保ってボックス(x,y,w,h)内に収める。"""
    from PIL import Image
    iw, ih = Image.open(path).size
    box_r = w / h; img_r = iw / ih
    if img_r > box_r:
        nw = w; nh = w / img_r
    else:
        nh = h; nw = h * img_r
    if align == "center":
        nx = x + (w - nw) / 2
    elif align == "left":
        nx = x
    else:
        nx = x + (w - nw)
    if valign == "middle":
        ny = y + (h - nh) / 2
    elif valign == "top":
        ny = y
    else:
        ny = y + (h - nh)
    return slide.shapes.add_picture(path, Inches(nx), Inches(ny), Inches(nw), Inches(nh))


def notes(slide, text_str):
    slide.notes_slide.notes_text_frame.text = text_str


def section_divider(prs, num, title, sub=""):
    """章扉：濃紺背景に白抜き。"""
    s = add_slide(prs, bg=INK_PANEL)
    rect(s, 0, 0, 0.28, 7.5, fill=ACCENT)
    text(s, 1.2, 2.75, 1.5, 1.0, [(f"{num:02d}", 44, ACCENT, FONT, {"bold": True})])
    text(s, 1.25, 3.55, 10.5, 1.2, [(title, 33, CARD, FONT, {"bold": True})])
    if sub:
        text(s, 1.3, 4.5, 10.3, 0.9, [(sub, 15, RGBColor(0xC7, 0xCE, 0xD4), FONT_L, {})])
    return s


def bubble(slide, x, y, w, h, role, role_color, role_bg, head, body_segs,
           answer=None, answer_ok=None, body_size=13.5):
    """会話の吹き出し。role=左の役割ラベル、body_segs=段落リスト、answer=右上の答えバッジ。
    上段（役割ラベル・見出し・答え）は固定、本文は上段の下からカード下端まで。"""
    card = rect(slide, x, y, w, h, fill=CARD, line=LINE, line_w=1.0,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05, shadow=True)
    # 左の色帯
    rect(slide, x, y, 0.12, h, fill=role_color, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    # 役割ラベル（上段左）
    rect(slide, x + 0.28, y + 0.16, 1.42, 0.4, fill=role_bg,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    text(slide, x + 0.28, y + 0.16, 1.42, 0.4, [(role, 12, role_color, FONT_SB, {})],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 見出し（上段中央）
    head_w = w - 3.2 if answer is not None else w - 2.0
    text(slide, x + 1.85, y + 0.16, head_w, 0.4, [(head, 11.5, SUB, FONT, {})],
         anchor=MSO_ANCHOR.MIDDLE)
    # 答えバッジ（上段右）
    if answer is not None:
        bcol = POS if answer_ok else NEG
        rect(slide, x + w - 1.28, y + 0.15, 1.04, 0.42, fill=bcol,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
        mark = "○" if answer_ok else "×"
        text(slide, x + w - 1.28, y + 0.15, 1.04, 0.42,
             [(f"{answer} {mark}", 12.5, CARD, FONT_SB, {})],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 本文（上段の下〜カード下端）
    para_block(slide, x + 0.34, y + 0.62, w - 0.66, h - 0.7, body_segs,
               line_spacing=1.14, space_after=3)
    return card
